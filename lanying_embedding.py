import lanying_redis
import logging
import uuid
import time
import re
import tiktoken
from markdownify import MarkdownConverter
import requests
import os
import random
import numpy as np
from redis.commands.search.query import Query
import pandas as pd
from pdfminer.high_level import extract_text
import hashlib
import subprocess
import docx2txt
from langchain.text_splitter import RecursiveCharacterTextSplitter
import lanying_url_loader
import json
import pdfplumber
import lanying_config
import lanying_pgvector
from urllib.parse import urlparse
from openai_token_counter import openai_token_counter
import lanying_chatbot
import lanying_config
import lanying_ai_capsule
import lanying_vendor
from pptx import Presentation

global_embedding_rate_limit = int(os.getenv("EMBEDDING_RATE_LIMIT", "30"))
global_embedding_lanying_connector_server = os.getenv("EMBEDDING_LANYING_CONNECTOR_SERVER", "https://lanying-connector.lanyingim.com")

tokenizer = tiktoken.get_encoding("cl100k_base")

class IgnoringScriptConverter(MarkdownConverter):
    """
    Create a custom MarkdownConverter that ignores script tags
    """
    def convert_script(self, el, text, convert_as_inline):
        return ''
    def convert_style(self, el, text, convert_as_inline):
        return ''


# Create shorthand method for conversion
def md(html, **options):
    return IgnoringScriptConverter(**options).convert(html)

def create_embedding(app_id, embedding_name, max_block_size, algo, admin_user_ids, preset_name, overlapping_size, vendor, model, type='text'):
    db_type = get_embedding_default_db_type(app_id)
    logging.info(f"start create embedding: app_id:{app_id}, embedding_name:{embedding_name}, max_block_size:{max_block_size},algo:{algo},admin_user_ids:{admin_user_ids},preset_name:{preset_name}, vendor:{vendor}, embedding_db_type: {db_type}")
    if app_id is None:
        app_id = ""
    old_embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if old_embedding_name_info:
        return {'result':"error", 'message': 'embedding_name exist'}
    model_config = lanying_vendor.get_embedding_model_config(vendor, model)
    if model_config is None:
        return {'result':"error", 'message': 'model not exist'}
    model = model_config['model']
    model_dim = model_config['dim']
    now = int(time.time())
    redis = lanying_redis.get_redis_stack_connection()
    embedding_uuid = generate_embedding_id()
    index_key = get_embedding_index_key(embedding_uuid)
    data_prefix_key = get_embedding_data_prefix_key(embedding_uuid)
    db_table_name = f"embedding_{embedding_uuid}_{app_id}"
    redis.hmset(get_embedding_name_key(app_id, embedding_name), {
        "app_id":app_id,
        "embedding_name": embedding_name,
        "embedding_uuid": embedding_uuid,
        "type": type,
        "time": now,
        "status": "ok",
        "admin_user_ids": ",".join([str(admin_user_id) for admin_user_id in admin_user_ids]),
        "preset_name":preset_name,
        "embedding_max_tokens":8192 if type == 'function' else 2048,
        "embedding_max_blocks":5,
        "embedding_content": "请严格按照下面的知识回答我之后的所有问题:"
    })
    if type != 'function':
        redis.rpush(get_embedding_names_key(app_id), embedding_name)
    redis.hmset(get_embedding_uuid_key(embedding_uuid),
                {"app_id": app_id,
                 "embedding_name": embedding_name,
                "index": index_key,
                "prefix": data_prefix_key,
                "max_block_size": max_block_size,
                "overlapping_size":overlapping_size,
                "algo": algo,
                "size": 0,
                "doc_id_seq": 0,
                "embedding_count":0,
                "embedding_size": 0,
                "vendor": vendor,
                "model": model,
                "text_size": 0,
                "time": now,
                "type": type,
                "db_type": db_type,
                "db_table_name": db_table_name,
                "status": "ok"})
    if db_type == 'redis':
        result = redis.execute_command("FT.CREATE", index_key, "prefix", "1", data_prefix_key, "SCHEMA","text","TEXT", "doc_id", "TAG", "embedding","VECTOR", "HNSW", "6", "TYPE", "FLOAT64","DIM", f"{model_dim}", "DISTANCE_METRIC",algo)
        logging.info(f"create_embedding success: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid} ft.create.result{result}")
    elif db_type == 'pgvector':
        with lanying_pgvector.get_connection() as conn:
            cursor = conn.cursor()
            cursor.execute(f"CREATE TABLE {db_table_name} (id bigserial PRIMARY KEY, embedding vector({model_dim}), content text, doc_id varchar(100),num_of_tokens int, summary text,text_hash varchar(100),question text,function text, reference text, block_id varchar(100));")
            cursor.execute(f"CREATE INDEX {db_table_name}_index_doc_id ON {db_table_name} (doc_id);")
            cursor.execute(f"CREATE INDEX {db_table_name}_index_embedding ON {db_table_name} USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100);")
            conn.commit()
            cursor.close()
            lanying_pgvector.put_connection(conn)
    update_app_embedding_admin_users(app_id, admin_user_ids)
    bind_preset_name(app_id, preset_name, embedding_name)
    return {'result':'ok', 'embedding_uuid':embedding_uuid}

def re_create_embedding_table(embedding_uuid):
    embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
    if embedding_uuid_info:
        db_type = embedding_uuid_info.get('db_type', 'redis')
        app_id = embedding_uuid_info['app_id']
        old_db_table_name = embedding_uuid_info.get('db_table_name', '')
        db_table_name = f"embedding_{embedding_uuid}_{app_id}_{int(time.time())}"
        vendor = embedding_uuid_info.get('vendor', 'openai')
        model = embedding_uuid_info.get('model', '')
        model_config = lanying_vendor.get_embedding_model_config(vendor, model)
        if model_config:
            model_dim = model_config['dim']
            if db_type == 'pgvector':
                with lanying_pgvector.get_connection() as conn:
                    cursor = conn.cursor()
                    cursor.execute(f"CREATE TABLE {db_table_name} (id bigserial PRIMARY KEY, embedding vector({model_dim}), content text, doc_id varchar(100),num_of_tokens int, summary text,text_hash varchar(100),question text,function text, reference text, block_id varchar(100));")
                    cursor.execute(f"CREATE INDEX {db_table_name}_index_doc_id ON {db_table_name} (doc_id);")
                    cursor.execute(f"CREATE INDEX {db_table_name}_index_embedding ON {db_table_name} USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100);")
                    conn.commit()
                    cursor.close()
                    lanying_pgvector.put_connection(conn)
                update_embedding_uuid_info(embedding_uuid, "db_table_name", db_table_name)
                redis = lanying_redis.get_redis_stack_connection()
                redis.rpush("lanying_connector:pgvector:table_to_deleted", old_db_table_name)
                return {'result': 'ok'}
    logging.error(f"fail to re_create_embedding_table: {embedding_uuid}")
    return {'result':'error'}

def migrate_embedding_from_redis_to_pgvector_one(app_id, embedding_name):
    redis = lanying_redis.get_redis_stack_connection()
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info is None:
        print(f"skip for not exist embedding_name: app_id:{app_id}, embedding_name:{embedding_name}")
        return
    embedding_uuid = embedding_name_info['embedding_uuid']
    embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
    if embedding_uuid_info is None:
        print(f"skip for not exist embedding_uuid: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}")
        return
    db_type = embedding_uuid_info.get('db_type', 'redis')
    if db_type == 'pgvector':
        print(f"skip for already pgvector: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}")
        return
    db_table_name = f"embedding_{embedding_uuid}_{app_id}"
    data_prefix = embedding_uuid_info.get('prefix', '')
    if data_prefix == '':
        print(f"skip for data_prefix is empty: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}")
        return
    with lanying_pgvector.get_connection() as conn:
        cursor = conn.cursor()
        cursor.execute(f"DROP TABLE IF EXISTS {db_table_name};")
        cursor.execute(f"CREATE TABLE {db_table_name} (id bigserial PRIMARY KEY, embedding vector(1536), content text, doc_id varchar(100),num_of_tokens int, summary text,text_hash varchar(100),question text,function text, reference text, block_id varchar(100));")
        cursor.execute(f"CREATE INDEX {db_table_name}_index_doc_id ON {db_table_name} (doc_id);")
        cursor.execute(f"CREATE INDEX {db_table_name}_index_embedding ON {db_table_name} USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100);")
        conn.commit()
        cursor.close()
        lanying_pgvector.put_connection(conn)
    key_count = 0
    for bytes in redis.scan_iter(match=f'{data_prefix}*', count=100):
        key_count += 1
        data_key = bytes.decode('utf-8')
        print(f"found data key: key_count:{key_count}, data_key:{data_key}")
        migrate_embedding_from_redis_to_pgvector_for_key(db_table_name, key_count, data_key, data_prefix, redis)
    print(f"update db_type: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}, db_table_name:{db_table_name}")
    redis.hmset(get_embedding_uuid_key(embedding_uuid),
                {"db_type": 'pgvector',
                "db_table_name": db_table_name})
    with lanying_pgvector.get_connection() as conn:
        cursor = conn.cursor()
        cursor.execute(f"SELECT COUNT(*) FROM {db_table_name}")
        row_count = cursor.fetchone()[0]
        cursor.close()
        lanying_pgvector.put_connection(conn)
        print(f"check row count: key_count:{key_count}, row_count:{row_count}, equal: {key_count == row_count}")
    print(f"finish migration: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}, db_table_name:{db_table_name}")

def redis_hgetall_with_embedding(redis, key):
    kvs = redis.hgetall(key)
    ret = {}
    if kvs:
        for k,v in kvs.items():
            new_k = k.decode('utf-8')
            if new_k == 'embedding':
                ret[new_k] = v
            else:
                ret[new_k] = v.decode('utf-8')
    return ret

def migrate_embedding_from_redis_to_pgvector_for_key(db_table_name, key_count, data_key, data_prefix, redis):
    info = redis_hgetall_with_embedding(redis, data_key)
    block_id = data_key[len(data_prefix):]
    text = info.get('text', '')
    question = info.get('question', '')
    text_hash = info.get('text_hash', '')
    function = info.get('function', '')
    embedding_bytes = info.get('embedding', b'')
    doc_id = info.get('doc_id', '')
    num_of_tokens = int(info.get('num_of_tokens', '0'))
    reference = info.get('reference', '')
    summary = info.get('summary', '{}')
    embedding = np.frombuffer(embedding_bytes, dtype=np.float64).tolist()
    if text_hash == '':
        embedding_text = text + question
        text_hash = sha256(embedding_text+function)
    with lanying_pgvector.get_connection() as conn:
        cursor = conn.cursor()
        insert_query = f"INSERT INTO {db_table_name} (embedding, content, doc_id, num_of_tokens, summary, text_hash, question, function, reference, block_id) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)"
        cursor.execute(insert_query, (embedding, text, doc_id, num_of_tokens, summary, text_hash, question, function, reference, block_id))
        conn.commit()
        cursor.close()
        lanying_pgvector.put_connection(conn)
    print(f"insert data finish | block_id:{block_id}, db_table_name:{db_table_name}, key_count:{key_count}, data_prefix:{data_prefix}")

def migrate_embedding_from_redis_to_pgvector_for_app_id(app_id):
    print(f"now processing app_id:{app_id}")
    embedding_names = list_embedding_names(app_id)
    for embedding_name in embedding_names:
        print(f"found embedding_name: {embedding_name}")
        migrate_embedding_from_redis_to_pgvector_one(app_id, embedding_name)

def migrate_embedding_from_redis_to_pgvector_all():
    redis = lanying_redis.get_redis_stack_connection()
    keys = lanying_redis.redis_keys(redis, "embedding_names:*")
    for key in keys:
        print(f"parsing key: {key}")
        fields = key.split(':')
        app_id = fields[1]
        print(f"found app_id: {app_id}")
        migrate_embedding_from_redis_to_pgvector_for_app_id(app_id)

def clean_embedding_from_redis_for_one(app_id, embedding_name):
    redis = lanying_redis.get_redis_stack_connection()
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info is None:
        print(f"skip for not exist embedding_name: app_id:{app_id}, embedding_name:{embedding_name}")
        return
    embedding_uuid = embedding_name_info['embedding_uuid']
    embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
    if embedding_uuid_info is None:
        print(f"skip for not exist embedding_uuid: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}")
        return
    db_type = embedding_uuid_info.get('db_type', 'redis')
    if db_type == 'redis':
        print(f"skip for redis db: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}")
        return
    db_table_name = f"embedding_{embedding_uuid}_{app_id}"
    data_prefix = embedding_uuid_info.get('prefix', '')
    if data_prefix == '':
        print(f"skip for data_prefix is empty: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}")
        return
    key_count = 0
    keys_to_delete = []
    for bytes in redis.scan_iter(match=f'{data_prefix}*', count=100):
        key_count += 1
        data_key = bytes.decode('utf-8')
        print(f"found data key: key_count:{key_count}, data_key:{data_key}")
        keys_to_delete.append(data_key)
    for key in keys_to_delete:
        if key.startswith('embedding_data:'):
            res = redis.delete(key)
            print(f"delete key:{key}, res:{res}")
        else:
            print(f"========================== skip delete for key:{data_key}")
    print(f"finish clean: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}, db_table_name:{db_table_name}")

def clean_embedding_from_redis_for_app_id(app_id):
    print(f"now processing app_id:{app_id}")
    embedding_names = list_embedding_names(app_id)
    for embedding_name in embedding_names:
        print(f"found embedding_name: {embedding_name}")
        clean_embedding_from_redis_for_one(app_id, embedding_name)

def clean_embedding_from_redis_all():
    redis = lanying_redis.get_redis_stack_connection()
    keys = lanying_redis.redis_keys(redis, "embedding_names:*")
    for key in keys:
        print(f"parsing key: {key}")
        fields = key.split(':')
        app_id = fields[1]
        print(f"found app_id: {app_id}")
        clean_embedding_from_redis_for_app_id(app_id)

def delete_embedding(app_id, embedding_name):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info:
        embedding_uuid = embedding_name_info["embedding_uuid"]
        embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
        if embedding_uuid_info:
            doc_id_list = get_embedding_doc_id_list(embedding_uuid, 0, -1)
            if len(doc_id_list) == 0:
                redis = lanying_redis.get_redis_stack_connection()
                redis.lrem(get_embedding_names_key(app_id), 1, embedding_name)
                redis.delete(get_embedding_uuid_key(embedding_uuid))
                redis.delete(get_embedding_name_key(app_id, embedding_name))
                return True
    return False

def configure_embedding(app_id, embedding_name, admin_user_ids, preset_name, embedding_max_tokens, embedding_max_blocks, embedding_content, new_embedding_name, max_block_size, overlapping_size, vendor, model):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info is None:
        return {'result':"error", 'message': 'embedding_name not exist'}
    if vendor == '':
        return {'result':"error", 'message': 'vendor not exist'}
    model_config = lanying_vendor.get_embedding_model_config(vendor, model)
    if model_config is None:
        return {'result':"error", 'message': 'model not exist'}
    model = model_config['model']
    embedding_uuid = embedding_name_info["embedding_uuid"]
    old_embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
    redis = lanying_redis.get_redis_stack_connection()
    if new_embedding_name != embedding_name and new_embedding_name != "":
        new_embedding_name_info = get_embedding_name_info(app_id, new_embedding_name)
        if new_embedding_name_info:
            logging.info(f"configure_embedding | new_embedding_name_info exists:{new_embedding_name_info}")
            return {'result':"error", 'message': 'new_embedding_name exist'}
        else:
            redis.rename(get_embedding_name_key(app_id,embedding_name), get_embedding_name_key(app_id,new_embedding_name))
            unbind_preset_name(app_id, embedding_name)
            list_key = get_embedding_names_key(app_id)
            redis.lrem(list_key, 1, embedding_name)
            redis.rpush(list_key, new_embedding_name)
            embedding_name = new_embedding_name
            embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
            if embedding_uuid_info:
                embedding_uuid_key = get_embedding_uuid_key(embedding_uuid)
                redis.hset(embedding_uuid_key, "embedding_name", new_embedding_name)
    redis.hmset(get_embedding_name_key(app_id, embedding_name), {
        "admin_user_ids": ",".join([str(admin_user_id) for admin_user_id in admin_user_ids]),
        "preset_name":preset_name,
        "embedding_name": embedding_name,
        "embedding_max_tokens":embedding_max_tokens,
        "embedding_max_blocks":embedding_max_blocks,
        "embedding_content": embedding_content
    })
    if max_block_size > 0:
        update_embedding_uuid_info(embedding_name_info['embedding_uuid'],"max_block_size", max_block_size)
    update_embedding_uuid_info(embedding_name_info['embedding_uuid'],"overlapping_size", overlapping_size)
    update_embedding_uuid_info(embedding_name_info['embedding_uuid'],"vendor", vendor)
    update_embedding_uuid_info(embedding_name_info['embedding_uuid'],"model", model)
    update_app_embedding_admin_users(app_id, admin_user_ids)
    bind_preset_name(app_id, preset_name, embedding_name)
    new_embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
    is_table_changed = False
    if old_embedding_uuid_info['vendor'] != new_embedding_uuid_info['vendor'] or old_embedding_uuid_info['model'] != new_embedding_uuid_info['model']:
        logging.info(f"configure_embedding recreate db table: {old_embedding_uuid_info}, {new_embedding_uuid_info}")
        re_create_result = re_create_embedding_table(embedding_uuid)
        embedding_type = new_embedding_uuid_info.get('type', 'text')
        if re_create_result['result'] == 'ok' and embedding_type != 'function':
            logging.info(f"configure_embedding re_run_all_doc_to_embedding: app_id:{app_id}, embedding_uuid:{embedding_uuid}")
            re_run_all_doc_to_embedding(app_id, embedding_uuid)
        is_table_changed = True
    return {"result":"ok", "data":{"is_table_changed": is_table_changed}}

def re_run_all_doc_to_embedding(app_id, embedding_uuid):
    trace_id = create_trace_id()
    doc_ids = get_embedding_doc_id_list(embedding_uuid, 0, -1)
    config = lanying_config.get_lanying_connector(app_id)
    update_embedding_uuid_info(embedding_uuid, "openai_secret_key", config['access_token'])
    from lanying_tasks import re_run_doc_to_embedding_by_doc_ids
    re_run_doc_to_embedding_by_doc_ids.apply_async(args = [trace_id, app_id, embedding_uuid, doc_ids])

def list_embeddings(app_id):
    redis = lanying_redis.get_redis_stack_connection()
    list_key = get_embedding_names_key(app_id)
    embedding_names = redis_lrange(redis, list_key, 0, -1)
    result = []
    for embedding_name in embedding_names:
        embedding_info = get_embedding_info_with_details(app_id, embedding_name)
        if embedding_info:
            result.append(embedding_info)
    return result

def get_embedding_info_with_details(app_id, embedding_name):
    embedding_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_info:
        embedding_info['admin_user_ids'] = embedding_info['admin_user_ids'].split(',')
        embedding_uuid = embedding_info["embedding_uuid"]
        embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
        for key in ["max_block_size","algo","embedding_count","embedding_size","text_size", "token_cnt", "preset_name", "embedding_max_tokens", "embedding_max_blocks", "embedding_content", "char_cnt", "storage_file_size", "overlapping_size", "vendor", "model"]:
            if key in embedding_uuid_info:
                embedding_info[key] = embedding_uuid_info[key]
        if "embedding_content" not in embedding_info:
            embedding_info["embedding_content"] = "请严格按照下面的知识回答我之后的所有问题:"
        return embedding_info
    
def list_embedding_names(app_id):
    redis = lanying_redis.get_redis_stack_connection()
    list_key = get_embedding_names_key(app_id)
    return redis_lrange(redis, list_key, 0, -1)

def search_embeddings(app_id, embedding_name, doc_id, embedding, max_tokens, max_blocks, is_fulldoc, doc_ids, check_storage_limit = True):
    if max_blocks > 100:
        max_blocks = 100
    result = []
    for extra_block_num in [10, 100, 500, 2000, 5000]:
        page_size = max_blocks+extra_block_num
        is_finish,result = search_embeddings_internal(app_id, embedding_name, doc_id, embedding, max_tokens, max_blocks, is_fulldoc, page_size, doc_ids, check_storage_limit)
        if is_finish:
            break
        logging.info(f"search_embeddings not finish: app_id:{app_id}, page_size:{page_size}, extra_block_num:{extra_block_num}")
    return result

def search_in_redis(app_id, embedding_name, doc_id, embedding, max_tokens, max_blocks, is_fulldoc, page_size, embedding_index, doc_ids):
    redis = lanying_redis.get_redis_stack_connection()
    if len(doc_ids) > 0:
        base_query = query_by_doc_ids(doc_ids)
    elif doc_id == "":
        base_query = f"*=>[KNN {page_size} @embedding $vector AS vector_score]"
    elif is_fulldoc:
        base_query = query_by_doc_id(doc_id)
    else:
        base_query = f"{query_by_doc_id(doc_id)}=>[KNN {page_size} @embedding  $vector AS vector_score]"
    query = Query(base_query).return_fields("text", "vector_score", "num_of_tokens", "summary","doc_id","text_hash","question","function", "reference").paging(0,page_size).dialect(2)
    if not is_fulldoc:
        query = query.sort_by("vector_score")
    return redis.ft(embedding_index).search(query, query_params={"vector": np.array(embedding).tobytes()})

def show_blocks(app_id, embedding_name, doc_id, count):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info:
        embedding_uuid = embedding_name_info['embedding_uuid']
        embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
        db_table_name = embedding_uuid_info['db_table_name']
        with lanying_pgvector.get_connection() as conn:
            cursor = conn.cursor()
            query = f"SELECT id,content,doc_id,num_of_tokens,summary,text_hash,question,function,reference,block_id FROM {db_table_name} where doc_id = %s ORDER BY block_id LIMIT %s;"
            args =  [doc_id, count]
            cursor.execute(query, args)
            rows = cursor.fetchall()
            cursor.close()
            lanying_pgvector.put_connection(conn)
            names = ['id','text','doc_id','num_of_tokens','summary','text_hash','question','function','reference','block_id']
            ret = []
            for row in rows:
                data = {}
                for index,name in enumerate(names):
                    data[name] = row[index]
                ret.append(data)
            return ret

def search_in_pgvector(app_id, embedding_name, doc_id, embedding, max_tokens, max_blocks, is_fulldoc, page_size, embedding_uuid_info, doc_ids):
    page_size = int(page_size)
    db_table_name = embedding_uuid_info['db_table_name']
    db_ivfflat_probes = int(embedding_uuid_info.get('db_ivfflat_probes', '32'))
    embedding_count = int(embedding_uuid_info.get('embedding_count', '0'))
    max_embedding_count = int(os.getenv('MAX_EMBEDDING_COUNT_FOR_FULL_TABLE_SCAN', '5000'))
    if embedding_count < max_embedding_count and db_ivfflat_probes < 100:
        logging.info(f"using full table scan: embedding_count:{embedding_count}")
        db_ivfflat_probes = 100
    start_time = time.time()
    with lanying_pgvector.get_connection() as conn:
        cursor = conn.cursor()
        embedding_str = f"{embedding}"
        if len(doc_ids) > 0:
            query = f"SELECT id,content,doc_id,num_of_tokens,summary,text_hash,question,function,reference,block_id,embedding <=> %s AS vector_score FROM {db_table_name} where doc_id in %s ORDER BY embedding <=> %s LIMIT %s;"
            args =  [embedding_str, tuple(doc_ids), embedding_str, page_size]
        elif doc_id == "":
            query = f"SELECT id,content,doc_id,num_of_tokens,summary,text_hash,question,function,reference,block_id,embedding <=> %s AS vector_score FROM {db_table_name} ORDER BY embedding <=> %s LIMIT %s;"
            args = [embedding_str, embedding_str, page_size]
        elif is_fulldoc:
            query = f"SELECT id,content,doc_id,num_of_tokens,summary,text_hash,question,function,reference,block_id,'0.0' AS vector_score FROM {db_table_name} where doc_id = %s ORDER BY doc_id LIMIT %s;"
            args =  [doc_id, page_size]
        else:
            query = f"SELECT id,content,doc_id,num_of_tokens,summary,text_hash,question,function,reference,block_id,embedding <=> %s AS vector_score FROM {db_table_name} where doc_id = %s ORDER BY embedding <=> %s LIMIT %s;"
            args = [embedding_str, doc_id, embedding_str, page_size]
        # logging.info(f"query:{query},args:{args}")
        cursor.execute(f"SET LOCAL ivfflat.probes = {db_ivfflat_probes};")
        cursor.execute(query, args)
        rows = cursor.fetchall()
        cursor.close()
        lanying_pgvector.put_connection(conn)
        # logging.info(f"rows:{rows}")
        logging.info(f"query finish with time: {time.time() - start_time}, db_table_name:{db_table_name}")
        class MyDocument:
            pass
        results = MyDocument()
        docs = []
        names = ['id','text','doc_id','num_of_tokens','summary','text_hash','question','function','reference','block_id', 'vector_score']
        for row in rows:
            doc = MyDocument()
            for index,name in enumerate(names):
                doc.__dict__[name] = row[index]
            logging.info(f"doc: block_id:{doc.block_id}, num_of_tokens:{doc.num_of_tokens}, text:\n{doc.text}, function:\n{doc.function}")
            docs.append(doc)
        setattr(results, 'docs', docs)
        return results

def search_embeddings_internal(app_id, embedding_name, doc_id, embedding, max_tokens, max_blocks, is_fulldoc, page_size, doc_ids, check_storage_limit):
    if check_storage_limit:
        result = check_storage_size(app_id)
        if result['result'] == 'error':
            logging.info(f"search_embeddings | skip search for exceed storage limit app, app_id:{app_id}, embedding_name:{embedding_name}")
            return (True, [])
    logging.info(f"search_embeddings_internal | app_id:{app_id}, embedding_name:{embedding_name}, doc_id:{doc_id}, max_tokens:{max_tokens}, max_blocks:{max_blocks}, is_fulldoc:{is_fulldoc}, page_size:{page_size}")
    redis = lanying_redis.get_redis_stack_connection()
    if redis:
        embedding_index = get_embedding_index(app_id, embedding_name)
        if embedding_index:
            embedding_name_info = get_embedding_name_info(app_id, embedding_name)
            embedding_uuid_info = get_embedding_uuid_info(embedding_name_info['embedding_uuid'])
            db_type = embedding_uuid_info.get('db_type', 'redis')
            if db_type == 'pgvector':
                results = search_in_pgvector(app_id, embedding_name, doc_id, embedding, max_tokens, max_blocks, is_fulldoc, page_size, embedding_uuid_info, doc_ids)
            else:
                results = search_in_redis(app_id, embedding_name, doc_id, embedding, max_tokens, max_blocks, is_fulldoc, page_size, embedding_index, doc_ids)
            # logging.info(f"topk result:{results.docs[:1]}")
            ret = []
            now_tokens = 0
            blocks_num = 0
            is_finish = False
            text_hashes = {'-'}
            docs = results.docs
            if is_fulldoc:
                docs_for_sort = []
                index = 0
                for doc in docs:
                    index = index + 1
                    seg_id = parse_segment_id_int_value(doc)
                    docs_for_sort.append(((seg_id, index),doc))
                docs = []
                for _,doc in sorted(docs_for_sort):
                    docs.append(doc)
            if len(docs) < page_size:
                logging.info(f"search_embeddings finish for no more doc: doc_count:{len(docs)}, page_size:{page_size}")
                is_finish = True
            max_continue_cnt = 5
            for doc in docs:
                text_hash = doc.text_hash if hasattr(doc, 'text_hash') else sha256(doc.text)
                if text_hash in text_hashes:
                    continue
                text_hashes.add(text_hash)
                now_tokens += int(doc.num_of_tokens)
                blocks_num += 1
                logging.info(f"search_embeddings count token: max_tokens:{max_tokens}, now_tokens:{now_tokens}, num_of_tokens:{int(doc.num_of_tokens)},blocks_num:{blocks_num}")
                if now_tokens > max_tokens:
                    if max_continue_cnt > 0:
                        max_continue_cnt -= 1
                        now_tokens -= int(doc.num_of_tokens)
                        is_finish = True
                        logging.info(f"search_embeddings num_of_token too large so skip: num_of_tokens:{int(doc.num_of_tokens)}, max_continue_cnt:{max_continue_cnt}")
                        continue
                    else:
                        is_finish = True
                        break
                if blocks_num > max_blocks:
                    is_finish = True
                    break
                ret.append(doc)
            return (is_finish, ret)
    return (True, [])

def get_preset_embedding_infos(embeddings, app_id, preset_name):
    if embeddings is None:
        redis = lanying_redis.get_redis_stack_connection()
        key = get_preset_name_key(app_id)
        bind_infos = redis_hgetall(redis, key)
        embedding_infos = []
        for now_embedding_name, now_preset_name in bind_infos.items():
            if now_preset_name == preset_name:
                embedding_info = get_embedding_name_info(app_id, now_embedding_name)
                if embedding_info:
                    embedding_info["embedding_name"] = now_embedding_name
                    embedding_info["embedding_uuid"] = embedding_info["embedding_uuid"]
                    embedding_info["embedding_max_tokens"] = int(embedding_info.get("embedding_max_tokens","2048"))
                    embedding_info["embedding_max_blocks"] = int(embedding_info.get("embedding_max_blocks", "5"))
                    embedding_infos.append(embedding_info)
        return embedding_infos
    else:
        embedding_uuids = embeddings.get(preset_name, [])
        embedding_infos = []
        for embedding_uuid in embedding_uuids:
            embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
            if embedding_uuid_info:
                embedding_info = get_embedding_name_info(app_id, embedding_uuid_info["embedding_name"])
                if embedding_info:
                    embedding_info["embedding_name"] = embedding_uuid_info["embedding_name"]
                    embedding_info["embedding_uuid"] = embedding_info["embedding_uuid"]
                    embedding_info["embedding_max_tokens"] = int(embedding_info.get("embedding_max_tokens","2048"))
                    embedding_info["embedding_max_blocks"] = int(embedding_info.get("embedding_max_blocks", "5"))
                    embedding_infos.append(embedding_info)
        return embedding_infos

def get_preset_embedding_infos_by_publish_capsule_id(capsule_id):
    capsule = lanying_ai_capsule.get_publish_capsule(capsule_id)
    if capsule:
        capsule_app_id = capsule['app_id']
        capsule_chatbot_id = capsule['chatbot_id']
        chatbot = lanying_chatbot.get_chatbot(capsule_app_id, capsule_chatbot_id)
        if chatbot:
            config = lanying_config.get_lanying_connector(capsule_app_id)
            return get_preset_embedding_infos(config.get('embeddings'), capsule_app_id, chatbot['name'])
    return []

def get_preset_embedding_infos_by_capsule_id(capsule_id):
    capsule = lanying_ai_capsule.get_capsule(capsule_id)
    if capsule:
        capsule_app_id = capsule['app_id']
        capsule_chatbot_id = capsule['chatbot_id']
        chatbot = lanying_chatbot.get_chatbot(capsule_app_id, capsule_chatbot_id)
        if chatbot:
            config = lanying_config.get_lanying_connector(capsule_app_id)
            return get_preset_embedding_infos(config.get('embeddings'), capsule_app_id, chatbot['name'])
    return []

def get_embedding_index(app_id, embedding_name):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info:
        embedding_uuid = embedding_name_info["embedding_uuid"]
        embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
        if embedding_uuid_info:
            return embedding_uuid_info["index"]
    global_embedding_name_info = get_global_embedding_name_info(embedding_name)
    if global_embedding_name_info:
        return global_embedding_name_info["index"]
    return None

def get_global_embedding_name_info(embedding_name):
    redis = lanying_redis.get_redis_stack_connection()
    global_alias_key = get_global_embedding_name_key(embedding_name)
    info = redis_hgetall(redis, global_alias_key)
    if "index" in info:
        return info

def get_global_embedding_name_key(embedding_name):
    return f"embedding_config:alias:{embedding_name}"

def process_embedding_file(trace_id, app_id, embedding_uuid, filename, origin_filename, doc_id, ext):
    redis = lanying_redis.get_redis_stack_connection()
    increase_embedding_doc_field(redis, embedding_uuid, doc_id, "process_count", 1)
    embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
    # logging.info(f"process_embedding_file | config:{embedding_uuid_info}")
    if embedding_uuid_info:
        try:
            if ext in [".html", ".htm"]:
                process_html(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
            elif ext in [".csv"]:
                process_csv(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
            elif ext in [".txt"]:
                process_txt(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
            elif ext in [".pdf"]:
                process_pdf(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
            elif ext in [".md"]:
                process_markdown(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
            elif ext in [".docx", ".doc"]:
                process_docx(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
            elif ext in [".xlsx", ".xls"]:
                process_xlsx(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
            elif ext in [".pptx"]:
                process_pptx(embedding_uuid_info, app_id, embedding_uuid, filename, origin_filename, doc_id)
        except Exception as e:
            increase_embedding_doc_field(redis, embedding_uuid, doc_id, "fail_count", 1)
            raise e
        increase_embedding_doc_field(redis, embedding_uuid, doc_id, "succ_count", 1)
        update_doc_field(embedding_uuid, doc_id, "status", "finish")
    else:
        logging.info(f"process_embedding_file embedding_uuid not exist | embedding_uuid:{embedding_uuid}")

def generate_embedding_id():
    redis = lanying_redis.get_redis_stack_connection()
    return redis.incrby("embedding_id_generator", 1)

def get_embedding_name_info(app_id, embedding_name):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_name_key(app_id, embedding_name)
    info = redis_hgetall(redis, key)
    if "embedding_uuid" in info:
        return info
    return None

def get_embedding_uuid_info(embedding_uuid):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_uuid_key(embedding_uuid)
    info = redis_hgetall(redis, key)
    if "index" in info:
        if 'vendor' not in info:
            vendor = 'openai'
            info['vendor'] = vendor
            update_embedding_uuid_info(embedding_uuid, "vendor", vendor)
        if 'model' not in info or info['model'] == '':
            vendor = info['vendor']
            model_config = lanying_vendor.get_embedding_model_config(vendor, '')
            if model_config:
                model = model_config['model']
                info['model'] = model
                update_embedding_uuid_info(embedding_uuid, "model", model)
            else:
                info['model'] = ''
        return info
    return None

def get_app_embedding_uuid_info(app_id, embedding_uuid):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_uuid_key(embedding_uuid)
    info = redis_hgetall(redis, key)
    if "app_id" in info and info['app_id'] == app_id:
        return info
    return None

def update_embedding_uuid_info(embedding_uuid, field, value):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_uuid_key(embedding_uuid)
    redis.hset(key, field, value)

def remove_space_line(text):
    lines = text.split('\n')
    new_lines = [line for line in lines if not re.match(r'^\s*$', line)]
    return '\n'.join(new_lines)

def estimate_html(embedding_uuid, html):
    redis = lanying_redis.get_redis_stack_connection()
    embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
    blocks = markdown_to_blocks(embedding_uuid_info, md(html))
    return len(blocks)

def process_html(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    with open(filename, "r") as f:
        html = f.read()
        process_markdown_content(config, app_id, embedding_uuid, origin_filename, doc_id, md(html))

def process_markdown_content(config, app_id, embedding_uuid, origin_filename, doc_id, markdown):
    blocks = markdown_to_blocks(config, markdown)
    redis = lanying_redis.get_redis_stack_connection()
    update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), len(blocks))
    insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, blocks, redis)

def markdown_to_blocks(config, markdown):
    markdown = remove_space_line(markdown)
    rule = config.get('block_split_rule',"^#{1,3} ")
    blocks = []
    total_tokens = 0
    for block in re.split(rule ,markdown, flags=re.MULTILINE):
        block_tokens, block_blocks = process_block(config, block)
        total_tokens += block_tokens
        blocks.extend(block_blocks)
    return blocks

def process_markdown(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()
        process_markdown_content(config, app_id, embedding_uuid, origin_filename, doc_id, content)

def process_txt(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    blocks = []
    total_tokens = 0
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()
        total_tokens, blocks = process_block(config, content)
    redis = lanying_redis.get_redis_stack_connection()
    update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), len(blocks))
    insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, blocks, redis)


def process_pdf(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    blocks = []
    total_tokens = 0
    content = extract_pdf(filename)
    total_tokens, blocks = process_block(config, content)
    redis = lanying_redis.get_redis_stack_connection()
    update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), len(blocks))
    insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, blocks, redis)

def extract_pdf(filename):
    try:
        with pdfplumber.open(filename) as pdf:
            texts = []
            tables = []
            for page in pdf.pages:
                texts.append(page.extract_text())
                for table in page.extract_tables():
                    rows = []
                    for row in table:
                        rows.append(f"{row}")
                    tables.append("\n".join(rows))
            return "\n".join(texts) + "\n\n" + "\n\n".join(tables)
    except Exception as e:
        logging.info("failed to extract pdf by pdfplumber")
        logging.exception(e)
        return extract_text(filename)

def process_docx(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    blocks = []
    total_tokens = 0
    ext = parse_file_ext(origin_filename)
    text = ''
    if ext == ".doc":
        try:
            output = subprocess.check_output(['antiword', filename])
            text = output.decode('utf-8', 'ignore')
        except subprocess.CalledProcessError:
            logging.error("Failed to convert the document: app_id:{app_id}, filename:{filename}, doc_id:{doc_id}")
            raise
    else:
        text = docx2txt.process(filename)
    total_tokens, blocks = process_block(config, text)
    redis = lanying_redis.get_redis_stack_connection()
    update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), len(blocks))
    insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, blocks, redis)

def process_pptx(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    logging.info(f"start process_pptx: embedding_uuid={embedding_uuid}, filename={filename}, origin_filename:{origin_filename}")
    redis = lanying_redis.get_redis_stack_connection()
    all_texts = extract_text_from_pptx(filename)
    max_block_size = get_max_token_count(config)
    now_texts = []
    now_tokens = 0
    threshold = 10
    total_tokens = 0
    total_blocks = 0
    for text in all_texts:
        tokens = num_of_tokens(text)
        if tokens + now_tokens + threshold > max_block_size:
            now_texts_str = '\n'.join(now_texts)
            block_tokens, block_blocks = process_block(config, now_texts_str)
            total_tokens += block_tokens
            total_blocks += len(block_blocks)
            update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
            insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)
            now_texts = [text]
            now_tokens = tokens
        else:
            now_texts.append(text)
            now_tokens += tokens
    if now_tokens > 0:
        now_texts_str = '\n'.join(now_texts)
        block_tokens, block_blocks = process_block(config, now_texts_str)
        total_tokens += block_tokens
        total_blocks += len(block_blocks)
        update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
        insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)

def extract_text_from_pptx(pptx_path):
    # 加载PPTX文件
    presentation = Presentation(pptx_path)
    all_texts = []

    # 遍历每一张幻灯片
    for slide in presentation.slides:
        slide_texts = []
        
        # 遍历每一个形状
        for shape in slide.shapes:
            if shape.has_text_frame:
                # 提取形状中的文本
                text = ''
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text
                    text += '\n'
                slide_texts.append(text)
        
        all_texts.append('\n'.join(slide_texts))

    return all_texts

def process_csv(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    logging.info(f"start process_csv: embedding_uuid={embedding_uuid}, filename={filename}")
    df = pd.read_csv(filename)
    df = df.fillna('')
    size = len(df)
    logging.info(f"embeddings: size={size}")
    total_blocks = 0
    total_tokens = 0
    redis = lanying_redis.get_redis_stack_connection()
    columns = df.columns.to_list()
    if 'text' in columns or ('question' in columns and 'answer' in columns):
        for i, row in df.iterrows():
            if 'question' in row and 'answer' in row:
                block_tokens, block_blocks = process_question(config, str(row['question']), str(row['answer']), row.get('reference',''))
                total_tokens += block_tokens
                total_blocks += len(block_blocks)
                update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
                insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)
            # elif 'text' in row and 'function' in row:
            #     block_tokens, block_blocks = process_function(config, str(row['text']), str(row['function']))
            #     total_tokens += block_tokens
            #     total_blocks += len(block_blocks)
            #     update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
            #     insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)
            elif 'text' in row:
                block_tokens, block_blocks = process_block(config, str(row['text']))
                total_tokens += block_tokens
                total_blocks += len(block_blocks)
                update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
                insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)
    else:
        df = pd.read_csv(filename, header=None)
        lines = []
        for i, row in df.iterrows():
            line_blocks = []
            for text in row:
                line_blocks.append(str(text))
            line = '\t'.join(line_blocks)
            lines.append(line)
        content = '\n'.join(lines)
        total_tokens, blocks = process_block(config, content)
        update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), len(blocks))
        insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, blocks, redis)

def process_xlsx(config, app_id, embedding_uuid, filename, origin_filename, doc_id):
    logging.info(f"start process_xlsx: embedding_uuid={embedding_uuid}, filename={filename}")
    xl_file = pd.ExcelFile(filename)
    total_blocks = 0
    total_tokens = 0
    for sheet_name in xl_file.sheet_names:
        df = pd.read_excel(filename, sheet_name=sheet_name)
        df = df.fillna('')
        size = len(df)
        logging.info(f"embeddings: size={size}, sheet_name={sheet_name}")
        redis = lanying_redis.get_redis_stack_connection()
        columns = df.columns.to_list()
        if 'text' in columns or ('question' in columns and 'answer' in columns):
            for i, row in df.iterrows():
                if 'question' in row and 'answer' in row:
                    block_tokens, block_blocks = process_question(config, str(row['question']), str(row['answer']), row.get('reference',''))
                    total_tokens += block_tokens
                    total_blocks += len(block_blocks)
                    update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
                    insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)
                # elif 'text' in row and 'function' in row:
                #     block_tokens, block_blocks = process_function(config, str(row['text']), str(row['function']))
                #     total_tokens += block_tokens
                #     total_blocks += len(block_blocks)
                #     update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
                #     insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)
                elif 'text' in row:
                    block_tokens, block_blocks = process_block(config, str(row['text']))
                    total_tokens += block_tokens
                    total_blocks += len(block_blocks)
                    update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
                    insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)
        else:
            df = pd.read_excel(filename, sheet_name=sheet_name, header=None)
            lines = []
            for i, row in df.iterrows():
                line_blocks = []
                for text in row:
                    line_blocks.append(str(text))
                line = '\t'.join(line_blocks)
                lines.append(line)
            content = '\n'.join(lines)
            block_tokens, block_blocks = process_block(config, content)
            total_tokens += block_tokens
            total_blocks += len(block_blocks)
            redis = lanying_redis.get_redis_stack_connection()
            update_progress_total(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), total_blocks)
            insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, block_blocks, redis)

def insert_embeddings(config, app_id, embedding_uuid, origin_filename, doc_id, blocks, redis):
    vendor = config.get('vendor', 'openai')
    advised_model = config.get('model', '')
    model_config = lanying_vendor.get_embedding_model_config(vendor, advised_model)
    model = model_config['model']
    db_type = config.get('db_type', 'redis')
    is_dry_run = config.get("dry_run", "false") == "true"
    max_block_size = get_max_token_count(config)
    question_answer_index_mode = config.get("question_answer_index_mode", "all")
    logging.info(f"insert_embeddings | app_id:{app_id}, embedding_uuid:{embedding_uuid}, origin_filename:{origin_filename}, doc_id:{doc_id}, is_dry_run:{is_dry_run}, block_count:{len(blocks)}, dry_run_from_config:{config.get('dry_run', 'None')}, vendor:{vendor}, model:{model}, max_block_size:{max_block_size}")
    for block in blocks:
        if len(block) == 2:
            token_cnt,text = block
            question = ''
            function = ''
            reference = ''
            advised_block_id = ''
        elif len(block) == 5 and block[1] == 'question':
            token_cnt, _, question, text, reference = block
            function = ''
            advised_block_id = ''
        elif len(block) == 5 and block[1] == 'function':
            token_cnt, _, text, function, advised_block_id = block
            question = ''
            reference = ''
        else:
            raise Exception(f"bad_block: {block}")
        doc_info = get_doc(embedding_uuid, doc_id)
        if doc_info:
            block_id = advised_block_id if len(advised_block_id) > 0 else generate_block_id(embedding_uuid, doc_id)
            maybe_rate_limit(5)
            if question_answer_index_mode == "question":
                if len(question) > 0:
                    embedding_text = question
                else:
                    embedding_text = question + text
            else:
                embedding_text = question + text
            embedding = fetch_embedding(app_id, vendor, model_config, embedding_text, is_dry_run)
            key = get_embedding_data_key(embedding_uuid, block_id)
            embedding_bytes = np.array(embedding).tobytes()
            text_hash = sha256(embedding_text+function)
            if db_type == "pgvector":
                db_table_name = config['db_table_name']
                def insert_fun():
                    with lanying_pgvector.get_connection() as conn:
                        cursor = conn.cursor()
                        insert_query = f"INSERT INTO {db_table_name} (embedding, content, doc_id, num_of_tokens, summary, text_hash, question, function, reference, block_id) VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s)"
                        cursor.execute(insert_query, (embedding, text, doc_id, token_cnt, "{}", text_hash, question, function, reference, block_id))
                        conn.commit()
                        cursor.close()
                        lanying_pgvector.put_connection(conn)
                retry_time = 10
                for i in range(retry_time):
                    try:
                        insert_fun()
                        break
                    except Exception as e:
                        logging.exception(e)
                        if i == retry_time -1:
                            logging.info(f"insert embedding fail at last retry:{i}")
                            raise e
                        else:
                            logging.info(f"insert embedding fail, schedule retry:{i}")
                        time.sleep(2)
            else:
                redis.hmset(key, {"text":text,
                                "question": question,
                                "text_hash":text_hash,
                                "embedding":embedding_bytes,
                                "doc_id": doc_id,
                                "num_of_tokens": token_cnt,
                                "function": function,
                                "reference": reference,
                                "summary": "{}"})
            embedding_size = len(embedding_bytes)
            text_size = text_byte_size(embedding_text)
            char_cnt = len(embedding_text)
            increase_embedding_uuid_field(redis, embedding_uuid, "embedding_count", 1)
            increase_embedding_uuid_field(redis, embedding_uuid, "embedding_size", embedding_size)
            increase_embedding_uuid_field(redis, embedding_uuid, "text_size", text_size)
            increase_embedding_uuid_field(redis, embedding_uuid, "char_cnt", char_cnt)
            increase_embedding_uuid_field(redis, embedding_uuid, "token_cnt", token_cnt)
            increase_embedding_doc_field(redis, embedding_uuid, doc_id, "embedding_count", 1)
            increase_embedding_doc_field(redis, embedding_uuid, doc_id, "embedding_size", embedding_size)
            increase_embedding_doc_field(redis, embedding_uuid, doc_id, "text_size", text_size)
            increase_embedding_doc_field(redis, embedding_uuid, doc_id, "char_cnt", char_cnt)
            increase_embedding_doc_field(redis, embedding_uuid, doc_id, "token_cnt", token_cnt)
            update_doc_field(embedding_uuid, doc_id, "vendor", vendor)
            update_doc_field(embedding_uuid, doc_id, "model", model)
            update_progress(redis, get_embedding_doc_info_key(embedding_uuid, doc_id), 1)
            question_desc = ''
            if question != '':
                question_desc = f"question:{question}\nanswer:"
            logging.info(f"=======block_id:{block_id},token_cnt:{token_cnt},char_cnt:{char_cnt},text_size:{text_size},max_block_size:{max_block_size}, text_hash:{text_hash}=====\n{question_desc}{text}")

def fetch_embedding(app_id, vendor, model_config, text, is_dry_run=False, retry = 10, sleep = 0.2, sleep_multi=1.7):
    model = model_config['model']
    if is_dry_run:
        return [random.uniform(0, 1) for i in range(model_config['dim'])]
    auth_secret = lanying_config.get_embedding_auth_secret()
    headers = {"Content-Type": "application/json",
               "Authorization": f"Bearer {app_id}-{auth_secret}"}
    body = {
        "text": text,
        "vendor": vendor,
        "model": model
    }
    response = {}
    url = global_embedding_lanying_connector_server + "/fetch_embeddings"
    try:
        response = requests.post(url, headers=headers, json = body).json()
        return response['embedding']
    except Exception as e:
        logging.info(f"fetch_embedding got error response:{response}, text:{text}, vendor:{vendor}, model:{model}")
        code = ""
        try:
            code = response["code"]
        except Exception as ee:
            pass
        if code in ["bad_authorization","no_quota", "deduct_failed"]:
            raise Exception(code)
        if retry > 0:
            time.sleep(sleep)
            return fetch_embedding(app_id, vendor, model_config, text, is_dry_run, retry-1, sleep * sleep_multi, sleep_multi)
        raise e

def process_block(config, block):
    block = remove_space_line(block)
    block = block.replace('\0','')
    chunks = []
    total_tokens = 0
    max_block_size = get_max_token_count(config)
    overlapping_size = get_overlapping_size(config)
    text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
            encoding_name="cl100k_base",
            chunk_size=max_block_size,
            chunk_overlap=overlapping_size,
            disallowed_special=(),
            separators=["\n", "。",".", "，", ","," ", ""]
        )
    texts = text_splitter.split_text(block)
    for text in texts:
        text = text.lstrip(",.，。 \n")
        if len(text) > 0:
            token_count = num_of_tokens(text)
            total_tokens += token_count
            chunks.append((token_count,text))
    return (total_tokens, chunks)

def process_question(config, question, answer, reference):
    question_token_cnt = num_of_tokens(question)
    answer_token_cnt = num_of_tokens(answer)
    token_limit = embedding_model_token_limit()
    token_cnt = question_token_cnt + answer_token_cnt
    if question_token_cnt == 0 or answer_token_cnt == 0:
        return (0,[])
    if token_cnt <= token_limit:
        return (token_cnt, [(token_cnt, "question", question, answer, reference)])
    else:
        logging.info(f"process_question | skip too large question answer: question_token_cnt:{question_token_cnt}, answer_token_cnt:{answer_token_cnt}")
        return (0, [])

# def process_function(config, text, function):
#     try:
#         function_obj = json.loads(function)
#         if "name" in function_obj and "description" in function_obj and "parameters" in function_obj and 'function_call' in function_obj:
#             text_token_cnt = num_of_tokens(text)
#             function_token_cnt = num_of_tokens(function)
#             token_limit = embedding_model_token_limit()
#             token_cnt = function_token_cnt
#             if text_token_cnt == 0:
#                 return (0,[])
#             if text_token_cnt <= token_limit:
#                 return (token_cnt, [(token_cnt, "function", text, function)])
#             else:
#                 logging.info(f"process_function | skip too large function: text_token_cnt:{text_token_cnt}, function_token_cnt:{function_token_cnt}")
#                 return (0, [])
#     except Exception as e:
#         logging.exception(e)
#         logging.info(f"fail to process_function | text:{text}, function:{function}")
#     return {0, []}

def embedding_model_token_limit():
    return 8000

def process_line(text, max_block_size):
    result = []
    current_str = ""
    token_cnt = 0
    char_cnt = 0
    for char in text:
        if token_cnt >= max_block_size:
            result.append((token_cnt, current_str))
            current_str = ""
            token_cnt = 0
            char_cnt = 0
        current_str += char
        char_cnt += 1
        if token_cnt >= max_block_size - 20 or char_cnt % 20 == 0:
            token_cnt = num_of_tokens(current_str)
    if len(current_str) > 0:
        token_cnt = num_of_tokens(current_str)
        result.append((token_cnt, current_str))
    return result

def get_max_token_count(config):
    max_block_size = max(350, int(config.get('max_block_size', "350")))
    max_token_count = word_num_to_token_num(max_block_size)
    if config.get('vendor', 'openai') == 'baidu':
        return min(max_token_count, 300)
    return max_token_count

def get_overlapping_size(config):
    overlapping_size = max(0, int(config.get('overlapping_size', "0")))
    token_count = word_num_to_token_num(overlapping_size)
    return token_count

def word_num_to_token_num(word_num):
    return round(word_num * 1.3)

def generate_block_id(embedding_uuid, doc_id):
    key = get_embedding_doc_info_key(embedding_uuid,doc_id)
    redis = lanying_redis.get_redis_stack_connection()
    result = redis.hincrby(key, "block_id_seq", 1)
    return f"{doc_id}-{result}"

def num_of_tokens(str):
    return len(tokenizer.encode(str, disallowed_special=()))

def maybe_rate_limit(retry):
    redis = lanying_redis.get_redis_stack_connection()
    now = time.time()
    key = f"embedding:rate_limit:{int(now)}"
    count = redis.incrby(key, 1)
    if count == 1:
        redis.expire(key, 10)
    if count > global_embedding_rate_limit and retry > 0:
        time.sleep(int(now)+1 - now + 0.1 * random.random())
        maybe_rate_limit(retry-1)

def create_trace_id():
    redis = lanying_redis.get_redis_stack_connection()
    for i in range(100):
        trace_id = str(uuid.uuid4())
        key = trace_id_key(trace_id)
        if redis.hsetnx(key, "status", "wait") > 0:
            return trace_id

def trace_finish(trace_id, app_id, status, message, doc_id, embedding_name):
    notify_user = get_trace_field(trace_id, "notify_user")
    if notify_user:
        auth_secret = lanying_config.get_embedding_auth_secret()
        headers = {"Content-Type": "application/json",
                "Authorization": f"Bearer {app_id}-{auth_secret}"}
        body = {
            "trace_id": trace_id,
            "status":status,
            "message": message,
            "doc_id":doc_id,
            "embedding_name": embedding_name
        }
        response = {}
        url = global_embedding_lanying_connector_server + "/trace_finish"
        try:
            response = requests.post(url, headers=headers, json = body)
            logging.info(f"trace_finish got response:{response}")
        except Exception as e:
            logging.exception(e)

def update_trace_field(trace_id, field, value):
    redis = lanying_redis.get_redis_stack_connection()
    key = trace_id_key(trace_id)
    redis.hset(key, field, value)

def get_all_trace_fields(trace_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = trace_id_key(trace_id)
    return redis_hgetall(redis, key)

def get_trace_field(trace_id, field):
    redis = lanying_redis.get_redis_stack_connection()
    key = trace_id_key(trace_id)
    return redis_hget(redis, key, field)

def delete_trace_field(trace_id, field):
    redis = lanying_redis.get_redis_stack_connection()
    key = trace_id_key(trace_id)
    redis.hdel(key, field)

def add_trace_doc_id(trace_id, doc_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = trace_doc_key(trace_id)
    redis.rpush(key, doc_id)

def clear_trace_doc_id(trace_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = trace_doc_key(trace_id)
    redis.delete(key)

def get_trace_doc_ids(trace_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = trace_doc_key(trace_id)
    return redis_lrange(redis, key, 0, -1)

def trace_id_key(trace_id):
    return f"lanying_trace_id:{trace_id}"

def trace_doc_key(trace_id):
    return f"lanying_trace_doc:{trace_id}"

def update_progress(redis, key, value):
    redis.hincrby(key, "progress_finish", value)

def update_progress_total(redis, key, total):
    redis.hset(key, "progress_total", total)

def increase_embedding_uuid_field(redis, embedding_uuid, field, value):
    redis.hincrby(get_embedding_uuid_key(embedding_uuid), field, value)

def update_app_embedding_admin_users(app_id, admin_user_ids):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_app_embedding_admin_user_key(app_id)
    for user_id in admin_user_ids:
        redis.hincrby(key, user_id, 1)

def bind_preset_name(app_id, preset_name, embedding_name):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_preset_name_key(app_id)
    redis.hset(key, embedding_name, preset_name)

def unbind_preset_name(app_id, embedding_name):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_preset_name_key(app_id)
    redis.hdel(key, embedding_name)

def is_app_embedding_admin_user(app_id, user_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_app_embedding_admin_user_key(app_id)
    return redis.hget(key, user_id) is not None

def get_embedding_names_key(app_id):
    return f"embedding_names:{app_id}"

def get_preset_name_key(app_id):
    return f"preset_names:{app_id}"

def get_embedding_name_key(app_id, embedding_name):
    return f"embedding_name:{app_id}:{embedding_name}"

def get_embedding_uuid_key(embedding_uuid):
    return f"embedding_uuid:{embedding_uuid}"

def get_embedding_index_key(embedding_uuid):
    return f"embedding_index:{embedding_uuid}"

def get_embedding_data_prefix_key(embedding_uuid):
    return f"embedding_data:{embedding_uuid}:"

def get_embedding_data_key(embedding_uuid, block_id):
    return f"embedding_data:{embedding_uuid}:{block_id}"

def get_app_embedding_admin_user_key(app_id):
    return f"embedding_admin_user_id:{app_id}"

def create_task(embedding_uuid, type, urls):
    redis = lanying_redis.get_redis_stack_connection()
    embedding_uuid_key = get_embedding_uuid_key(embedding_uuid)
    task_id_seq = redis.hincrby(embedding_uuid_key, "task_id_seq", 1)
    task_id = f"{embedding_uuid}-{task_id_seq}"
    info_key = get_embedding_task_info_key(embedding_uuid, task_id)
    task_list_key = get_embedding_task_list_key(embedding_uuid)
    redis.hmset(info_key, {
        "embedding_uuid":embedding_uuid,
        "type": type,
        "url": urls[0],
        "time": int(time.time()),
        "status": "wait",
        "visited_num": 0,
        "to_visit_num": 0,
        "found_num": 0,
        "file_size": 0,
        "block_num": 0,
        "processing_total_num": 0,
        "processing_success_num": 0,
        "processing_fail_num": 0
    })
    redis.rpush(task_list_key, task_id)
    return task_id

def get_task_list(embedding_uuid):
    task_list = []
    redis = lanying_redis.get_redis_stack_connection()
    task_list_key = get_embedding_task_list_key(embedding_uuid)
    for task_id in redis_lrange(redis, task_list_key, 0, -1):
        task_info = get_task(embedding_uuid, task_id)
        if task_info:
            if 'site_task_id' in task_info:
                site_task_id = task_info['site_task_id']
                site_task_info = lanying_url_loader.info_task(site_task_id)
                task_info['visited_num'] = site_task_info['visited']
                task_info['to_visit_num'] = site_task_info['to_visit']
                task_info['found_num'] = site_task_info['found']
            task_info["task_id"] = task_id
            task_list.append(task_info)
    return task_list

def delete_task(embedding_uuid, task_id):
    task_info = get_task(embedding_uuid, task_id)
    if task_info:
        redis = lanying_redis.get_redis_stack_connection()
        info_key = get_embedding_task_info_key(embedding_uuid, task_id)
        task_list_key = get_embedding_task_list_key(embedding_uuid)
        task_detail_key = get_embedding_task_detail_key(embedding_uuid, task_id)
        redis.delete(info_key, task_detail_key)
        redis.lrem(task_list_key,2, task_id)
        if 'site_task_id' in task_info:
            lanying_url_loader.clean_task(task_info['site_task_id'])

def get_task(embedding_uuid, task_id):
    redis = lanying_redis.get_redis_stack_connection()
    info_key = get_embedding_task_info_key(embedding_uuid, task_id)
    info = redis_hgetall(redis, info_key)
    if "type" in info:
        return info
    return None

def get_task_details(embedding_uuid, task_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_task_detail_key(embedding_uuid, task_id)
    return redis_hgetall(redis, key)

def get_task_details_iterator(embedding_uuid, task_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_task_detail_key(embedding_uuid, task_id)
    for k,v in redis.hscan_iter(key, '*'):
        json_data = json.loads(v)
        url = k.decode("utf-8")
        yield (url, json_data)

def get_task_details_count(embedding_uuid, task_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_task_detail_key(embedding_uuid, task_id)
    return redis.hlen(key)

def delete_task_details_by_fields(embedding_uuid, task_id, fields):
    if len(fields) > 0:
        redis = lanying_redis.get_redis_stack_connection()
        key = get_embedding_task_detail_key(embedding_uuid, task_id)
        redis.hdel(key, *fields)

def get_task_detail_fields(embedding_uuid, task_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_task_detail_key(embedding_uuid, task_id)
    return redis_hkeys(redis, key)

def get_task_detail_field(embedding_uuid, task_id, field):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_task_detail_key(embedding_uuid, task_id)
    return redis.hget(key, field)

def set_task_detail_field(embedding_uuid, task_id, field, value):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_embedding_task_detail_key(embedding_uuid, task_id)
    redis.hset(key, field, value)

def update_task_field(embedding_uuid, task_id, field, value):
    redis = lanying_redis.get_redis_stack_connection()
    info_key = get_embedding_task_info_key(embedding_uuid, task_id)
    return redis.hset(info_key, field, value)

def increase_task_field(embedding_uuid, task_id, field, value):
    redis = lanying_redis.get_redis_stack_connection()
    return redis.hincrby(get_embedding_task_info_key(embedding_uuid, task_id), field, value)

def create_doc_info(app_id, embedding_uuid, filename, object_name, doc_id, file_size, ext, type, source, vendor, opts):
    redis = lanying_redis.get_redis_stack_connection()
    info_key = get_embedding_doc_info_key(embedding_uuid, doc_id)
    lanying_link = ''
    if opts.get('generate_lanying_links', False) == True:
        if type in ['url', 'site'] and filename.startswith("http") and is_in_lanying_link_white_list(filename):
            link_res = generate_lanying_links(app_id, filename)
            if link_res['result'] == 'ok':
                lanying_link = link_res['link']
    logging.info(f"create_doc_info | app_id:{app_id}, embedding_uuid:{embedding_uuid}, filename:{filename}, doc_id:{doc_id}, type:{type}, lanying_link:{lanying_link}, opts:{opts}")
    redis.hmset(info_key, {"filename":filename,
                           "object_name":object_name,
                           "time": int(time.time()),
                           "file_size": file_size,
                           "ext": ext,
                           "type": type,
                           "source": source,
                           "vendor": vendor,
                           "lanying_link": lanying_link,
                           "status": "wait"})

def generate_lanying_links(app_id, long_link):
    authorization = os.getenv("LANYING_CONNECTOR_LANYING_LINK_API_KEY")
    if authorization:
        url = "https://lanying.link/api/generate_prefix_url"
        headers = {
            'authorization': authorization
        }
        body = {
            'prefix': 'doc/',
            'long_link': long_link,
            'app_id': app_id,
            'uid': 0
        }
        try:
            response = requests.post(url, headers= headers, json = body)
            response_json = response.json()
            if response_json.get('code') == 200:
                return {'result': 'ok', 'link' : response_json.get('data')}
            else:
                return {'result': 'error', 'reason':'bad_status_code'}
        except Exception as e:
            return {'result': 'error', 'reason':'exception'}
    else:
        return {'result': 'error', 'reason': 'no_authorization'}

def update_doc_field(embedding_uuid, doc_id, field, value):
    redis = lanying_redis.get_redis_stack_connection()
    info_key = get_embedding_doc_info_key(embedding_uuid, doc_id)
    redis.hset(info_key, field, value)

def add_doc_to_embedding(embedding_uuid, doc_id):
    redis = lanying_redis.get_redis_stack_connection()
    list_key = get_embedding_doc_list_key(embedding_uuid)
    redis.rpush(list_key, doc_id)

def delete_doc_from_embedding(app_id, embedding_name, doc_id, task):
    logging.info(f"delete_doc_from_embedding started | app_id:{app_id}, embedding_name:{embedding_name}, doc_id:{doc_id}")
    redis = lanying_redis.get_redis_stack_connection()
    embedding_index = get_embedding_index(app_id, embedding_name)
    if embedding_index:
        embedding_name_info = get_embedding_name_info(app_id, embedding_name)
        embedding_uuid = embedding_name_info["embedding_uuid"]
        doc_info = get_doc(embedding_uuid, doc_id)
        embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
        if doc_info and embedding_uuid_info:
            db_type = embedding_uuid_info.get('db_type', 'redis')
            db_table_name = embedding_uuid_info.get('db_table_name', '')
            task.apply_async(args = [app_id, embedding_name, doc_id, embedding_index, 0, db_type, db_table_name])
            list_key = get_embedding_doc_list_key(embedding_uuid)
            redis.lrem(list_key, 2, doc_id)
            restore_storage_size(app_id, embedding_uuid, doc_id)
            info_key = get_embedding_doc_info_key(embedding_uuid, doc_id)
            redis.delete(info_key)
            increase_embedding_uuid_field(redis, embedding_uuid, "embedding_count", -int(doc_info.get("embedding_count", "0")))
            increase_embedding_uuid_field(redis, embedding_uuid, "embedding_size", -int(doc_info.get("embedding_size", "0")))
            increase_embedding_uuid_field(redis, embedding_uuid, "text_size", -int(doc_info.get("text_size", "0")))
            increase_embedding_uuid_field(redis, embedding_uuid, "token_cnt", -int(doc_info.get("token_cnt", "0")))
            increase_embedding_uuid_field(redis, embedding_uuid, "char_cnt", -int(doc_info.get("char_cnt", "0")))
            return True
    return False

def search_doc_data_and_delete(app_id, embedding_name, doc_id, embedding_index, last_total, db_type, db_table_name):
    logging.info(f"delete doc_id from embedding| db_type={db_type}, app_id={app_id}, embedding_name={embedding_name}, doc_id={doc_id}")
    if db_type == 'pgvector':
        with lanying_pgvector.get_connection() as conn:
            cursor = conn.cursor()
            cursor.execute(f"delete from {db_table_name} where doc_id = %s", [doc_id])
            conn.commit()
            cursor.close()
            lanying_pgvector.put_connection(conn)
    else:
        redis = lanying_redis.get_redis_stack_connection()
        base_query = query_by_doc_id(doc_id)
        query = Query(base_query).no_content().paging(0, 200).dialect(2)
        results = redis.ft(embedding_index).search(query)
        logging.info(f"search for delete  | app_id:{app_id}, embedding_name:{embedding_name}, doc_id:{doc_id}, last_total:{last_total}, result:{results}")
        if results.total == 0 or results.total == last_total:
            logging.info(f"search for delete stop  | app_id:{app_id}, embedding_name:{embedding_name}, doc_id:{doc_id}, last_total:{last_total}")
            return
        keys = []
        for doc in results.docs:
            keys.append(doc.id)
        if len(keys) > 0:
            redis.delete(*keys)
        if len(keys) < results.total:
            search_doc_data_and_delete(app_id, embedding_name, doc_id, embedding_index, results.total, db_type, db_table_name)
        else:
            logging.info(f"search for delete stop for last page | app_id:{app_id}, embedding_name:{embedding_name}, doc_id:{doc_id}, last_total:{last_total}")

def delete_embedding_block(app_id, embedding_name, doc_id, block_id):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info:
        embedding_uuid = embedding_name_info["embedding_uuid"]
        embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
        db_type = embedding_uuid_info.get('db_type', 'redis')
        if db_type == 'pgvector':
            with lanying_pgvector.get_connection() as conn:
                db_table_name = embedding_uuid_info['db_table_name']
                cursor = conn.cursor()
                cursor.execute(f"delete from {db_table_name} where doc_id = %s and block_id = %s", [doc_id, block_id])
                conn.commit()
                cursor.close()
                lanying_pgvector.put_connection(conn)
        else:
            redis = lanying_redis.get_redis_stack_connection()
            key = get_embedding_data_key(embedding_uuid, block_id)
            redis.delete(key)

def query_by_doc_id(doc_id):
    if len(doc_id) < 30:
        new_doc_id = doc_id.replace('-','\\-')
        return "@doc_id:{"+new_doc_id+"}"
    else: # for deprecated doc_id format
        return "@doc_id:{"+doc_id+"}"

def query_by_doc_ids(doc_ids):
    doc_ids_str = ','.join(doc_ids)
    new_doc_ids_str = doc_ids_str.replace('-','\\-')
    return "@doc_id:{"+new_doc_ids_str+"}"

def get_doc(embedding_uuid, doc_id):
    redis = lanying_redis.get_redis_stack_connection()
    info_key = get_embedding_doc_info_key(embedding_uuid, doc_id)
    info = redis_hgetall(redis, info_key)
    if "filename" in info:
        return info
    return None

def get_embedding_name_by_doc_id(app_id, doc_id):
    fields = doc_id.split('-')
    if len(fields) == 2:
        embedding_uuid = fields[0]
        embedding_uuid_info = get_embedding_uuid_info(embedding_uuid)
        if embedding_uuid_info and embedding_uuid_info['app_id'] == app_id:
            return embedding_uuid_info['embedding_name']
    return None

def get_doc_metadata(app_id, embedding_name, doc_id):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info is None:
        return {'result': 'error', 'message': 'embedding name not exist'}
    embedding_uuid = embedding_name_info['embedding_uuid']
    doc_info = get_doc(embedding_uuid, doc_id)
    if doc_info is None:
        return {'result': 'error', 'message': 'doc_id not exist'}
    metadata = {}
    try:
        metadata = json.loads(doc_info.get('metadata', '{}'))
    except Exception as e:
        pass
    return {'result': 'ok', 'data': metadata}

def set_doc_metadata(app_id, embedding_name, doc_id, metadata):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    if embedding_name_info:
        embedding_uuid = embedding_name_info['embedding_uuid']
        doc_info = get_doc(embedding_uuid, doc_id)
        if doc_info:
            update_doc_field(embedding_uuid, doc_id, "metadata", json.dumps(metadata, ensure_ascii=False))

def check_storage_size(app_id):
    redis = lanying_redis.get_redis_stack_connection()
    storage_limit = get_app_config_int(app_id, "lanying_connector.storage_limit")
    storage_payg = get_app_config_int(app_id, "lanying_connector.storage_payg")
    if storage_payg == 1 and storage_limit > 0:
        return {'result':'ok'}
    now_storage_size = increase_app_storage_file_size(app_id, 0)
    if storage_limit == 0 or now_storage_size > storage_limit * 1024 * 1024 * 1.1:
        return {'result': 'error'}
    return {'result':'ok'}

def add_storage_size(app_id, embedding_uuid, doc_id, file_size):
    redis = lanying_redis.get_redis_stack_connection()
    storage_limit = get_app_config_int(app_id, "lanying_connector.storage_limit")
    storage_payg = get_app_config_int(app_id, "lanying_connector.storage_payg")
    if storage_payg == 1 and storage_limit > 0:
        increase_app_storage_file_size(app_id, file_size)
        increase_embedding_doc_field(redis, embedding_uuid, doc_id, "storage_file_size", file_size)
        increase_embedding_uuid_field(redis, embedding_uuid, "storage_file_size", file_size)
        return {'result':'ok'}
    now_storage_size = increase_app_storage_file_size(app_id, 0)
    if now_storage_size + file_size > storage_limit * 1024 * 1024:
        return {'result': 'error'}
    now_storage_size = increase_app_storage_file_size(app_id, file_size)
    if now_storage_size > storage_limit * 1024 * 1024:
        increase_app_storage_file_size(app_id, -file_size)
        return {'result': 'error'}
    increase_embedding_doc_field(redis, embedding_uuid, doc_id, "storage_file_size", file_size)
    increase_embedding_uuid_field(redis, embedding_uuid, "storage_file_size", file_size)
    return {'result':'ok'}

def restore_storage_size(app_id, embedding_uuid, doc_id):
    redis = lanying_redis.get_redis_stack_connection()
    doc_info = get_doc(embedding_uuid, doc_id)
    if doc_info:
        storage_file_size = int(doc_info.get('storage_file_size', "0"))
        if storage_file_size > 0:
            result = increase_embedding_doc_field(redis, embedding_uuid, doc_id, "storage_file_size", -storage_file_size)
            if result >= 0:
                increase_embedding_uuid_field(redis, embedding_uuid, "storage_file_size", -storage_file_size)
                increase_app_storage_file_size(app_id, -storage_file_size)

def increase_app_storage_file_size(app_id, value):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_app_embedding_app_info_key(app_id)
    result = redis.hincrby(key, "storage_file_size", value)
    storage_file_size_max = redis.hincrby(key, "storage_file_size_max", 0)
    if storage_file_size_max < result:
        redis.hset(key, "storage_file_size_max", result)
    return result

def increase_embedding_doc_field(redis, embedding_uuid, doc_id, field, value):
    return redis.hincrby(get_embedding_doc_info_key(embedding_uuid, doc_id), field, value)

def generate_doc_id(embedding_uuid):
    key = get_embedding_uuid_key(embedding_uuid)
    redis = lanying_redis.get_redis_stack_connection()
    doc_id_seq = redis.hincrby(key, "doc_id_seq", 1)
    return f"{embedding_uuid}-{doc_id_seq}"

def get_embedding_doc_info_list(app_id, embedding_name, start, end):
    embedding_name_info = get_embedding_name_info(app_id, embedding_name)
    total = 0
    result = []
    if embedding_name_info:
        embedding_uuid = embedding_name_info["embedding_uuid"]
        list_key = get_embedding_doc_list_key(embedding_uuid)
        redis = lanying_redis.get_redis_stack_connection()
        total = redis.llen(list_key)
        doc_id_list = redis_lrange(redis, list_key, start, end)
        for doc_id in doc_id_list:
            doc = get_doc(embedding_uuid,doc_id)
            if doc:
                doc['doc_id'] = doc_id
                logging.info(f"doc_info: app_id:{app_id}, embedding_name:{embedding_name}, embedding_uuid:{embedding_uuid}, doc_id:{doc_id}, info:{doc}")
                result.append(doc)
    return (total, result)

def get_embedding_doc_id_list(embedding_uuid, start, end):
    list_key = get_embedding_doc_list_key(embedding_uuid)
    redis = lanying_redis.get_redis_stack_connection()
    doc_id_list = redis_lrange(redis, list_key, start, end)
    return doc_id_list

def get_embedding_doc_list_key(embedding_uuid):
    return f"embedding_doc_list:{embedding_uuid}"

def get_embedding_doc_info_key(embedding_uuid, doc_id):
    return f"embedding_doc_info:{embedding_uuid}:{doc_id}"

def get_embedding_task_list_key(embedding_uuid):
    return f"embedding_task_list:{embedding_uuid}"

def get_embedding_task_info_key(embedding_uuid, task_id):
    return f"embedding_task_info:{embedding_uuid}:{task_id}"

def get_embedding_task_detail_key(embedding_uuid, task_id):
    return f"embedding_task_detail:{embedding_uuid}:{task_id}"

def get_app_embedding_app_info_key(app_id):
    return f"embedding_app_info:{app_id}"

def redis_lrange(redis, key, start, end):
    return [bytes.decode('utf-8') for bytes in redis.lrange(key, start, end)]

def redis_hkeys(redis, key):
    return [bytes.decode('utf-8') for bytes in redis.hkeys(key)]

def redis_hgetall(redis, key):
    kvs = redis.hgetall(key)
    ret = {}
    if kvs:
        for k,v in kvs.items():
            ret[k.decode('utf-8')] = v.decode('utf-8')
    return ret
def redis_hget(redis, key, field):
    result = redis.hget(key, field)
    if result:
        return result.decode('utf-8')
    return None

def text_byte_size(text):
    return len(text.encode('utf-8'))

def get_embedding_usage(app_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_app_embedding_app_info_key(app_id)
    return redis_hgetall(redis, key)

def set_embedding_usage(app_id, storage_file_size_max):
    redis = lanying_redis.get_redis_stack_connection()
    key = get_app_embedding_app_info_key(app_id)
    redis.hset(key, "storage_file_size_max", storage_file_size_max)
    return True

def save_app_config(app_id, key, value):
    if key.startswith("lanying_connector."):
        redis = lanying_redis.get_redis_stack_connection()
        name = get_app_config_key(app_id)
        redis.hset(name, key, value)

def get_app_config_int(app_id, key):
    redis = lanying_redis.get_redis_stack_connection()
    name = get_app_config_key(app_id)
    return redis.hincrby(name, key, 0)

def get_app_config_key(app_id):
    return lanying_config.get_redis_app_config_key(app_id)

def get_embedding_uuid_from_doc_id(doc_id):
    fields = doc_id.split('-')
    if len(fields) > 0:
        return fields[0]
    else:
        return None

def sha256(text):
    value = hashlib.sha256(text.encode('utf-8')).hexdigest()
    return value

def allow_exts():
    return [".html", ".htm", ".csv", ".txt", ".md", ".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx"]

def parse_file_ext(filename):
    if is_file_url(filename):
        return ".html"
    _,ext = os.path.splitext(filename)
    return ext.lower()

def is_file_url(filename):
    return filename.startswith("http://") or filename.startswith("https://")

def parse_segment_id_int_value(doc):
    if hasattr(doc, 'block_id'):
        seg_id = doc.block_id
    else:
        seg_id = doc.id
    fields = seg_id.split('-')
    try:
        return int(fields[len(fields)-1])
    except Exception as e:
        try:
            fields = seg_id.split(':')
            return int(fields[len(fields)-1])
        except Exception as ee:
            return 0

def get_embedding_default_db_type(app_id):
    redis = lanying_redis.get_redis_stack_connection()
    key = f"lanying_connector:embedding:default_db_type"
    value = lanying_redis.redis_hget(redis, key, '*')
    if value:
        return value
    value = lanying_redis.redis_hget(redis, key, app_id)
    if value:
        return value
    return 'redis'

def set_embedding_default_db_type(app_id, default_db_type):
    redis = lanying_redis.get_redis_stack_connection()
    key = f"lanying_connector:embedding:default_db_type"
    redis.hset(key, app_id, default_db_type)

def is_in_lanying_link_white_list(url):
    redis = lanying_redis.get_redis_stack_connection()
    parse_url= urlparse(url.strip(' '))
    domain = parse_url.netloc
    if len(domain) > 0:
        return redis.sismember("lanying_connector:lanying_link:whitelist", domain)
    return False

def get_preset_names(app_id):
    if lanying_chatbot.is_chatbot_mode(app_id):
        return lanying_chatbot.get_chatbot_names(app_id)
    else:
        preset_names = ["default"]
        config = lanying_config.get_lanying_connector(app_id)
        if "preset" in config and "presets" in config["preset"]:
            try:
                for k in config["preset"]["presets"].keys():
                    preset_names.append(k)
            except Exception as e:
                logging.exception(e)
                pass
        return preset_names

def calc_functions_tokens(functions, model, vendor):
    if len(functions) == 0:
        return 0
    try:
        if vendor == 'openai':
            if model in ["text-embedding-3-large", "text-embedding-3-small"]:
                token_cnt = openai_token_counter(messages=[], functions=functions, model='text-embedding-ada-002') - openai_token_counter(messages=[])
            else:
                token_cnt = openai_token_counter(messages=[], functions=functions, model=model) - openai_token_counter(messages=[])
        else:
            token_cnt = openai_token_counter(messages=[], functions=functions) - openai_token_counter(messages=[])
        return token_cnt
    except Exception as e:
        logging.exception(e)
        token_cnt = len(json.dumps(functions, ensure_ascii=False))
        logging.info(f"fallback token_cnt:{token_cnt}")
        return token_cnt

def calc_function_tokens(function, model, vendor):
    return calc_functions_tokens([function], model, vendor)
